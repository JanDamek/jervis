"""Document Extraction Service — converts any file format to plain text.

Supported formats:
- text/*, code files → direct UTF-8 decode
- CSV/TSV → structured table extraction
- HTML → BeautifulSoup strip tags
- DOCX → python-docx paragraphs + tables
- DOC (legacy) → mammoth → HTML → text
- XLSX → openpyxl sheet rows
- XLS (legacy) → xlrd sheet rows
- PPTX → python-pptx slides (shapes + tables + notes)
- ODT/ODS/ODP → odfpy text extraction
- RTF → striprtf
- PDF → pymupdf text + VLM for pages with images/scans (hybrid)
- Images → VLM description + OCR
- MSG (Outlook) → extract-msg
- EML → stdlib email parser
- EPUB → ebooklib chapters
- VSDX (Visio) → shape text extraction
- ZIP/RAR/7z/tar.gz → recursive extraction of contained files

VLM calls go through Ollama router (HTTP) for GPU/cloud routing.
"""

from __future__ import annotations

import base64
import io
import logging
import os
import zipfile
import tarfile
from dataclasses import dataclass, field
from pathlib import PurePosixPath

import httpx

logger = logging.getLogger(__name__)

# Max recursion depth for nested archives
_MAX_ARCHIVE_DEPTH = 3
# Max total extracted size from archive (100MB)
_MAX_ARCHIVE_TOTAL_BYTES = 100 * 1024 * 1024
# Max files to extract from a single archive
_MAX_ARCHIVE_FILES = 200
# Skip files larger than this inside archives (20MB)
_MAX_SINGLE_FILE_BYTES = 20 * 1024 * 1024


@dataclass
class ImageDescription:
    description: str
    ocr_text: str = ""


@dataclass
class PageContent:
    page_number: int
    text: str
    images: list[ImageDescription] = field(default_factory=list)


@dataclass
class ExtractedDocument:
    text: str
    pages: list[PageContent] = field(default_factory=list)
    language: str = ""
    method: str = "direct"
    metadata: dict = field(default_factory=dict)


# ═══════════════════════════════════════════════════════════════════════════════
# MIME type routing
# ═══════════════════════════════════════════════════════════════════════════════

_TEXT_MIMES = {
    "text/plain", "text/markdown", "text/x-markdown",
    "application/json", "application/xml", "text/xml", "text/yaml",
    "application/x-yaml",
}
_CSV_MIMES = {"text/csv", "text/tab-separated-values"}
_HTML_MIMES = {"text/html", "application/xhtml+xml"}
_PDF_MIMES = {"application/pdf"}
_DOCX_MIMES = {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
_DOC_MIMES = {"application/msword"}
_XLSX_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}
_XLS_MIMES = {"application/vnd.ms-excel"}
_PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
_ODT_MIMES = {"application/vnd.oasis.opendocument.text"}
_ODS_MIMES = {"application/vnd.oasis.opendocument.spreadsheet"}
_ODP_MIMES = {"application/vnd.oasis.opendocument.presentation"}
_RTF_MIMES = {"application/rtf", "text/rtf"}
_MSG_MIMES = {"application/vnd.ms-outlook"}
_EML_MIMES = {"message/rfc822"}
_EPUB_MIMES = {"application/epub+zip"}
_VSDX_MIMES = {"application/vnd.ms-visio.drawing.main+xml",
               "application/vnd.visio"}
_ZIP_MIMES = {"application/zip", "application/x-zip-compressed"}
_TAR_MIMES = {"application/x-tar", "application/gzip", "application/x-gzip",
              "application/x-bzip2", "application/x-xz",
              "application/x-compressed-tar"}
_RAR_MIMES = {"application/x-rar-compressed", "application/vnd.rar"}
_SEVENZ_MIMES = {"application/x-7z-compressed"}
_IMAGE_MIMES = {
    "image/png", "image/jpeg", "image/jpg", "image/webp",
    "image/bmp", "image/gif", "image/tiff",
}

_EXT_TO_TYPE = {
    # Text
    ".txt": "text", ".json": "text", ".xml": "text",
    ".yaml": "text", ".yml": "text", ".md": "text", ".log": "text",
    ".ini": "text", ".cfg": "text", ".conf": "text", ".properties": "text",
    ".toml": "text", ".env": "text",
    # CSV/TSV
    ".csv": "csv", ".tsv": "csv",
    # Code (as text)
    ".py": "text", ".js": "text", ".ts": "text", ".tsx": "text",
    ".jsx": "text", ".java": "text", ".kt": "text", ".kts": "text",
    ".swift": "text", ".go": "text", ".rs": "text", ".c": "text",
    ".cpp": "text", ".h": "text", ".hpp": "text", ".cs": "text",
    ".rb": "text", ".php": "text", ".sh": "text", ".bash": "text",
    ".zsh": "text", ".fish": "text", ".ps1": "text",
    ".sql": "text", ".gradle": "text", ".groovy": "text",
    ".scala": "text", ".clj": "text", ".ex": "text", ".exs": "text",
    ".r": "text", ".R": "text", ".m": "text", ".mm": "text",
    ".dart": "text", ".lua": "text", ".pl": "text", ".pm": "text",
    ".tf": "text", ".hcl": "text",  # Terraform
    ".proto": "text", ".graphql": "text", ".gql": "text",
    ".dockerfile": "text", ".makefile": "text",
    ".cmake": "text", ".bat": "text", ".cmd": "text",
    ".vue": "text", ".svelte": "text",
    ".css": "text", ".scss": "text", ".less": "text", ".sass": "text",
    # HTML
    ".html": "html", ".htm": "html", ".xhtml": "html",
    # PDF
    ".pdf": "pdf",
    # Office — modern
    ".docx": "docx", ".xlsx": "xlsx", ".pptx": "pptx",
    # Office — legacy
    ".doc": "doc", ".xls": "xls", ".ppt": "pptx",
    # OpenDocument
    ".odt": "odt", ".ods": "ods", ".odp": "odp",
    # RTF
    ".rtf": "rtf",
    # Email
    ".msg": "msg", ".eml": "eml", ".mbox": "eml",
    # EPUB
    ".epub": "epub",
    # Visio
    ".vsdx": "vsdx",
    # Archives
    ".zip": "zip", ".jar": "zip", ".war": "zip",
    ".tar": "tar", ".tar.gz": "tar", ".tgz": "tar",
    ".tar.bz2": "tar", ".tar.xz": "tar",
    ".gz": "tar",
    ".rar": "rar",
    ".7z": "7z",
    # Images
    ".png": "image", ".jpg": "image", ".jpeg": "image", ".webp": "image",
    ".bmp": "image", ".gif": "image", ".tiff": "image", ".tif": "image",
    ".svg": "html",  # SVG is XML — extract text
}


def _detect_type(mime_type: str, filename: str) -> str:
    mime_lower = mime_type.lower().strip() if mime_type else ""

    # MIME-based detection
    if mime_lower in _TEXT_MIMES: return "text"
    if mime_lower in _CSV_MIMES: return "csv"
    if mime_lower in _HTML_MIMES: return "html"
    if mime_lower in _PDF_MIMES: return "pdf"
    if mime_lower in _DOCX_MIMES: return "docx"
    if mime_lower in _DOC_MIMES: return "doc"
    if mime_lower in _PPTX_MIMES: return "pptx"
    if mime_lower in _ODT_MIMES: return "odt"
    if mime_lower in _ODS_MIMES: return "ods"
    if mime_lower in _ODP_MIMES: return "odp"
    if mime_lower in _RTF_MIMES: return "rtf"
    if mime_lower in _MSG_MIMES: return "msg"
    if mime_lower in _EML_MIMES: return "eml"
    if mime_lower in _EPUB_MIMES: return "epub"
    if mime_lower in _VSDX_MIMES: return "vsdx"
    if mime_lower in _ZIP_MIMES: return "zip"
    if mime_lower in _TAR_MIMES: return "tar"
    if mime_lower in _RAR_MIMES: return "rar"
    if mime_lower in _SEVENZ_MIMES: return "7z"
    if mime_lower in _IMAGE_MIMES: return "image"
    # xlsx vs xls — resolve by extension
    if mime_lower in _XLSX_MIMES or mime_lower in _XLS_MIMES:
        fname_lower = filename.lower() if filename else ""
        if fname_lower.endswith(".xls"):
            return "xls"
        return "xlsx"

    # Extension-based fallback (check longest match first for .tar.gz etc)
    fname_lower = filename.lower() if filename else ""
    for ext in sorted(_EXT_TO_TYPE.keys(), key=len, reverse=True):
        if fname_lower.endswith(ext):
            return _EXT_TO_TYPE[ext]

    # Unknown → try as text
    return "text"


# Binary file extensions — listed in archive output but content not extracted
_BINARY_EXTENSIONS = {
    # Executables / system
    ".exe", ".dll", ".so", ".dylib", ".lib", ".a", ".o", ".obj",
    ".bin", ".com", ".sys", ".drv", ".msi", ".dmg", ".app",
    # JVM
    ".class", ".jar", ".war", ".ear",
    # .NET
    ".pdb", ".nupkg",
    # Python
    ".pyc", ".pyo", ".pyd", ".whl", ".egg",
    # Native compiled
    ".deb", ".rpm", ".apk", ".ipa", ".aab",
    # Fonts
    ".ttf", ".otf", ".woff", ".woff2", ".eot",
    # Media (not processable as text or VLM-worthy in archive context)
    ".mp3", ".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv",
    ".wav", ".flac", ".aac", ".ogg", ".m4a",
    ".ico",
    # Database / data
    ".db", ".sqlite", ".sqlite3", ".mdb", ".accdb",
    # Compressed (nested archives handled by recursion, these are opaque)
    ".iso", ".img", ".vmdk", ".vhd",
    # Misc binary
    ".dat", ".pak", ".res", ".cache",
    ".min.js", ".min.css",  # minified — technically text but useless
    ".map",  # source maps — very large, not useful for KB
    ".lock",  # lock files
    ".wasm",
}


def _is_binary_extension(filename_lower: str) -> bool:
    """Check if file is a known binary format (no useful text to extract)."""
    for ext in _BINARY_EXTENSIONS:
        if filename_lower.endswith(ext):
            return True
    return False


def _human_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"


# ═══════════════════════════════════════════════════════════════════════════════
# Main extractor
# ═══════════════════════════════════════════════════════════════════════════════

class DocumentExtractor:
    """Document extraction service — any file → plain text."""

    def __init__(self, ollama_router_url: str = "http://jervis-ollama-router:8080"):
        self.router_url = ollama_router_url

    async def extract(
        self,
        file_bytes: bytes,
        filename: str,
        mime_type: str = "",
        max_tier: str = "NONE",
        _depth: int = 0,
    ) -> ExtractedDocument:
        doc_type = _detect_type(mime_type, filename)
        logger.info("DocumentExtractor: file=%s mime=%s type=%s size=%d depth=%d",
                     filename, mime_type, doc_type, len(file_bytes), _depth)

        # Sync extractors
        sync_extractors = {
            "text": lambda: self._extract_text_direct(file_bytes, filename),
            "csv": lambda: self._extract_csv(file_bytes, filename),
            "html": lambda: self._extract_html(file_bytes),
            "docx": lambda: self._extract_docx(file_bytes),
            "doc": lambda: self._extract_doc_legacy(file_bytes),
            "xlsx": lambda: self._extract_xlsx(file_bytes),
            "xls": lambda: self._extract_xls_legacy(file_bytes),
            "pptx": lambda: self._extract_pptx(file_bytes),
            "odt": lambda: self._extract_odf(file_bytes, "odt"),
            "ods": lambda: self._extract_odf(file_bytes, "ods"),
            "odp": lambda: self._extract_odf(file_bytes, "odp"),
            "rtf": lambda: self._extract_rtf(file_bytes),
            "msg": lambda: self._extract_msg(file_bytes),
            "eml": lambda: self._extract_eml(file_bytes),
            "epub": lambda: self._extract_epub(file_bytes),
            "vsdx": lambda: self._extract_vsdx(file_bytes),
        }

        if doc_type in sync_extractors:
            return sync_extractors[doc_type]()
        elif doc_type == "pdf":
            return await self._extract_pdf(file_bytes, max_tier)
        elif doc_type == "image":
            return await self._extract_image(file_bytes, max_tier)
        elif doc_type in ("zip", "tar", "rar", "7z"):
            return await self._extract_archive(file_bytes, filename, doc_type, max_tier, _depth)
        else:
            return self._extract_text_direct(file_bytes, filename)

    # ═══════════════════════════════════════════════════════════════════════════
    # Text / CSV / HTML
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_text_direct(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        text = file_bytes.decode("utf-8", errors="replace")
        return ExtractedDocument(text=text, method="direct", metadata={"filename": filename})

    def _extract_csv(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        """Extract CSV/TSV with proper table structure."""
        import csv
        text_content = file_bytes.decode("utf-8", errors="replace")

        # Detect delimiter
        try:
            dialect = csv.Sniffer().sniff(text_content[:4096])
            delimiter = dialect.delimiter
        except csv.Error:
            delimiter = "," if filename.lower().endswith(".csv") else "\t"

        reader = csv.reader(io.StringIO(text_content), delimiter=delimiter)
        rows = []
        header = None
        for i, row in enumerate(reader):
            if i == 0:
                header = row
                rows.append(" | ".join(row))
                rows.append(" | ".join(["---"] * len(row)))
            else:
                rows.append(" | ".join(row))
            if i > 5000:
                rows.append(f"... (truncated at {i} rows)")
                break

        text = "\n".join(rows)
        return ExtractedDocument(text=text, method="csv",
                                 metadata={"delimiter": delimiter, "columns": len(header) if header else 0})

    def _extract_html(self, file_bytes: bytes) -> ExtractedDocument:
        from bs4 import BeautifulSoup
        html_text = file_bytes.decode("utf-8", errors="replace")
        soup = BeautifulSoup(html_text, "lxml")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return ExtractedDocument(text=text, method="html",
                                 metadata={"title": soup.title.string if soup.title else ""})

    # ═══════════════════════════════════════════════════════════════════════════
    # Word: DOCX + DOC (legacy)
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_docx(self, file_bytes: bytes) -> ExtractedDocument:
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(file_bytes))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))
        text = "\n\n".join(parts)
        return ExtractedDocument(text=text, method="docx",
                                 metadata={"page_count": len(doc.sections), "has_tables": len(doc.tables) > 0})

    def _extract_doc_legacy(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            import mammoth
            result = mammoth.convert_to_html(io.BytesIO(file_bytes))
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(result.value, "lxml")
            text = soup.get_text(separator="\n", strip=True)
            return ExtractedDocument(text=text, method="doc-mammoth",
                                     metadata={"warnings": len(result.messages)})
        except Exception as e:
            logger.warning("Legacy .doc extraction failed: %s — trying as text", e)
            return self._extract_text_direct(file_bytes, "document.doc")

    # ═══════════════════════════════════════════════════════════════════════════
    # Excel: XLSX + XLS (legacy)
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_xlsx(self, file_bytes: bytes) -> ExtractedDocument:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    sheet_rows.append(" | ".join(cells))
            if sheet_rows:
                parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(sheet_rows))
        wb.close()
        text = "\n\n".join(parts)
        return ExtractedDocument(text=text, method="xlsx", metadata={"sheet_count": len(wb.sheetnames)})

    def _extract_xls_legacy(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=file_bytes)
            parts = []
            for sheet_name in wb.sheet_names():
                ws = wb.sheet_by_name(sheet_name)
                sheet_rows = []
                for row_idx in range(ws.nrows):
                    cells = [str(ws.cell_value(row_idx, col_idx)) for col_idx in range(ws.ncols)]
                    if any(c.strip() for c in cells):
                        sheet_rows.append(" | ".join(cells))
                if sheet_rows:
                    parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(sheet_rows))
            text = "\n\n".join(parts)
            return ExtractedDocument(text=text, method="xls-xlrd",
                                     metadata={"sheet_count": wb.nsheets})
        except Exception as e:
            logger.warning("Legacy .xls extraction failed: %s — trying openpyxl", e)
            return self._extract_xlsx(file_bytes)

    # ═══════════════════════════════════════════════════════════════════════════
    # PowerPoint: PPTX
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_pptx(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_parts.append(para.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_parts.append(" | ".join(cells))
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"[Notes: {notes}]")
                parts.append("\n".join(slide_parts))
            text = "\n\n".join(parts)
            return ExtractedDocument(text=text, method="pptx",
                                     metadata={"slide_count": len(prs.slides)})
        except Exception as e:
            logger.warning("PPTX extraction failed: %s — trying as text", e)
            return self._extract_text_direct(file_bytes, "presentation.pptx")

    # ═══════════════════════════════════════════════════════════════════════════
    # OpenDocument: ODT, ODS, ODP
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_odf(self, file_bytes: bytes, odf_type: str) -> ExtractedDocument:
        try:
            from odf.opendocument import load as odf_load
            from odf import text as odf_text, table as odf_table
            from odf.element import Element

            doc = odf_load(io.BytesIO(file_bytes))

            def get_text_recursive(element: Element) -> str:
                result = []
                if hasattr(element, 'childNodes'):
                    for child in element.childNodes:
                        if hasattr(child, 'data'):
                            result.append(child.data)
                        else:
                            result.append(get_text_recursive(child))
                return "".join(result)

            parts = []

            if odf_type == "odt":
                for para in doc.getElementsByType(odf_text.P):
                    t = get_text_recursive(para).strip()
                    if t:
                        parts.append(t)
                for heading in doc.getElementsByType(odf_text.H):
                    t = get_text_recursive(heading).strip()
                    if t:
                        parts.append(f"## {t}")
            elif odf_type in ("ods", "odp"):
                for tbl in doc.getElementsByType(odf_table.Table):
                    table_name = tbl.getAttribute("name") or ""
                    if table_name:
                        parts.append(f"=== {table_name} ===")
                    for row in tbl.getElementsByType(odf_table.TableRow):
                        cells = []
                        for cell in row.getElementsByType(odf_table.TableCell):
                            cells.append(get_text_recursive(cell).strip())
                        if any(cells):
                            parts.append(" | ".join(cells))
                for para in doc.getElementsByType(odf_text.P):
                    t = get_text_recursive(para).strip()
                    if t:
                        parts.append(t)

            text = "\n\n".join(parts) if parts else "(empty document)"
            return ExtractedDocument(text=text, method=f"odf-{odf_type}", metadata={"odf_type": odf_type})
        except Exception as e:
            logger.warning("ODF %s extraction failed: %s — trying as text", odf_type, e)
            return self._extract_text_direct(file_bytes, f"document.{odf_type}")

    # ═══════════════════════════════════════════════════════════════════════════
    # RTF
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_rtf(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            from striprtf.striprtf import rtf_to_text
            rtf_content = file_bytes.decode("utf-8", errors="replace")
            text = rtf_to_text(rtf_content)
            return ExtractedDocument(text=text, method="rtf")
        except Exception as e:
            logger.warning("RTF extraction failed: %s — trying as text", e)
            return self._extract_text_direct(file_bytes, "document.rtf")

    # ═══════════════════════════════════════════════════════════════════════════
    # Email: MSG (Outlook) + EML
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_msg(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            import extract_msg
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                msg = extract_msg.Message(tmp_path)
                parts = []
                if msg.subject:
                    parts.append(f"Subject: {msg.subject}")
                if msg.sender:
                    parts.append(f"From: {msg.sender}")
                if msg.to:
                    parts.append(f"To: {msg.to}")
                if msg.date:
                    parts.append(f"Date: {msg.date}")
                parts.append("")
                if msg.body:
                    parts.append(msg.body)
                text = "\n".join(parts)
                metadata = {
                    "subject": msg.subject or "",
                    "sender": msg.sender or "",
                    "attachment_count": len(msg.attachments),
                }
                msg.close()
                return ExtractedDocument(text=text, method="msg", metadata=metadata)
            finally:
                os.unlink(tmp_path)
        except Exception as e:
            logger.warning("MSG extraction failed: %s — trying as text", e)
            return self._extract_text_direct(file_bytes, "email.msg")

    def _extract_eml(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            import email
            from email import policy
            msg = email.message_from_bytes(file_bytes, policy=policy.default)
            parts = []
            if msg["subject"]:
                parts.append(f"Subject: {msg['subject']}")
            if msg["from"]:
                parts.append(f"From: {msg['from']}")
            if msg["to"]:
                parts.append(f"To: {msg['to']}")
            if msg["date"]:
                parts.append(f"Date: {msg['date']}")
            parts.append("")

            body = msg.get_body(preferencelist=('plain', 'html'))
            if body:
                content = body.get_content()
                if body.get_content_type() == "text/html":
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(content, "lxml")
                    content = soup.get_text(separator="\n", strip=True)
                parts.append(content)

            text = "\n".join(parts)
            return ExtractedDocument(text=text, method="eml",
                                     metadata={"subject": msg["subject"] or ""})
        except Exception as e:
            logger.warning("EML extraction failed: %s — trying as text", e)
            return self._extract_text_direct(file_bytes, "email.eml")

    # ═══════════════════════════════════════════════════════════════════════════
    # EPUB
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_epub(self, file_bytes: bytes) -> ExtractedDocument:
        try:
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(io.BytesIO(file_bytes))
            parts = []

            title = book.get_metadata('DC', 'title')
            if title:
                parts.append(f"# {title[0][0]}")

            for item in book.get_items_of_type(9):  # ITEM_DOCUMENT
                soup = BeautifulSoup(item.get_content(), "lxml")
                text = soup.get_text(separator="\n", strip=True)
                if text.strip():
                    parts.append(text)

            text = "\n\n".join(parts)
            return ExtractedDocument(text=text, method="epub",
                                     metadata={"title": title[0][0] if title else ""})
        except Exception as e:
            logger.warning("EPUB extraction failed: %s", e)
            return ExtractedDocument(text="(EPUB extraction failed)", method="epub-error")

    # ═══════════════════════════════════════════════════════════════════════════
    # Visio (VSDX)
    # ═══════════════════════════════════════════════════════════════════════════

    def _extract_vsdx(self, file_bytes: bytes) -> ExtractedDocument:
        """VSDX is a ZIP containing XML — extract shape text."""
        try:
            from bs4 import BeautifulSoup
            parts = []
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                for name in zf.namelist():
                    if name.startswith("visio/pages/page") and name.endswith(".xml"):
                        with zf.open(name) as f:
                            soup = BeautifulSoup(f.read(), "lxml-xml")
                            for text_elem in soup.find_all("Text"):
                                t = text_elem.get_text(strip=True)
                                if t:
                                    parts.append(t)
            text = "\n".join(parts) if parts else "(no text found in Visio diagram)"
            return ExtractedDocument(text=text, method="vsdx",
                                     metadata={"shape_count": len(parts)})
        except Exception as e:
            logger.warning("VSDX extraction failed: %s", e)
            return ExtractedDocument(text="(Visio extraction failed)", method="vsdx-error")

    # ═══════════════════════════════════════════════════════════════════════════
    # Archives: ZIP, TAR, RAR, 7z — recursive extraction
    # ═══════════════════════════════════════════════════════════════════════════

    async def _extract_archive(
        self, file_bytes: bytes, filename: str, archive_type: str,
        max_tier: str, depth: int,
    ) -> ExtractedDocument:
        if depth >= _MAX_ARCHIVE_DEPTH:
            return ExtractedDocument(
                text=f"(archive recursion limit reached: {filename})",
                method="archive-skip",
            )

        parts = []
        binary_listing = []
        files_processed = 0
        total_bytes = 0
        truncated = False

        try:
            entries = self._list_archive_entries(file_bytes, filename, archive_type)

            for entry_name, entry_bytes in entries:
                if files_processed >= _MAX_ARCHIVE_FILES:
                    truncated = True
                    break
                if total_bytes >= _MAX_ARCHIVE_TOTAL_BYTES:
                    truncated = True
                    break

                base = PurePosixPath(entry_name).name.lower()
                size = len(entry_bytes)

                # Skip OS metadata
                if base in ("thumbs.db", "desktop.ini", ".ds_store"):
                    continue

                # Binary files — list but don't extract content
                if _is_binary_extension(base):
                    binary_listing.append(f"  {entry_name} ({_human_size(size)})")
                    files_processed += 1
                    continue

                # Empty files
                if size == 0:
                    continue

                # Too large for extraction — list with size
                if size > _MAX_SINGLE_FILE_BYTES:
                    binary_listing.append(f"  {entry_name} ({_human_size(size)}, too large to extract)")
                    files_processed += 1
                    continue

                try:
                    result = await self.extract(
                        entry_bytes, entry_name, max_tier=max_tier, _depth=depth + 1,
                    )
                    if result.text.strip():
                        parts.append(f"=== {entry_name} ===\n{result.text}")
                        files_processed += 1
                        total_bytes += size
                except Exception as e:
                    logger.warning("Failed to extract %s from archive: %s", entry_name, e)
                    binary_listing.append(f"  {entry_name} ({_human_size(size)}, extraction error)")

        except Exception as e:
            logger.warning("Archive extraction failed for %s: %s", filename, e)
            return ExtractedDocument(
                text=f"(failed to open archive: {filename}: {e})",
                method="archive-error",
            )

        if binary_listing:
            parts.append("=== Binary/non-text files ===\n" + "\n".join(binary_listing))
        if truncated:
            parts.append(f"(archive truncated: limit reached at {files_processed} files)")

        text = "\n\n".join(parts) if parts else "(empty archive)"
        return ExtractedDocument(
            text=text, method=f"archive-{archive_type}",
            metadata={
                "archive_type": archive_type,
                "files_extracted": files_processed,
                "binary_files": len(binary_listing),
            },
        )

    def _list_archive_entries(
        self, file_bytes: bytes, filename: str, archive_type: str,
    ) -> list[tuple[str, bytes]]:
        """Yield (name, bytes) for each file in the archive."""
        entries = []

        if archive_type == "zip":
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                for info in zf.infolist():
                    if info.is_dir():
                        continue
                    try:
                        entries.append((info.filename, zf.read(info)))
                    except Exception as e:
                        logger.warning("Skipping %s in ZIP: %s", info.filename, e)

        elif archive_type == "tar":
            mode = "r:*"  # auto-detect compression
            try:
                with tarfile.open(fileobj=io.BytesIO(file_bytes), mode=mode) as tf:
                    for member in tf.getmembers():
                        if not member.isfile():
                            continue
                        f = tf.extractfile(member)
                        if f:
                            entries.append((member.name, f.read()))
            except Exception:
                # Try plain gzip single file
                import gzip
                try:
                    data = gzip.decompress(file_bytes)
                    inner_name = filename.replace(".gz", "")
                    entries.append((inner_name, data))
                except Exception:
                    pass

        elif archive_type == "rar":
            try:
                import rarfile
                with rarfile.RarFile(io.BytesIO(file_bytes)) as rf:
                    for info in rf.infolist():
                        if info.is_dir():
                            continue
                        entries.append((info.filename, rf.read(info)))
            except ImportError:
                logger.warning("RAR support requires 'rarfile' package + unrar binary")
            except Exception as e:
                logger.warning("RAR extraction failed: %s", e)

        elif archive_type == "7z":
            try:
                import py7zr
                with py7zr.SevenZipFile(io.BytesIO(file_bytes), mode='r') as sz:
                    for fname, bio in sz.readall().items():
                        entries.append((fname, bio.read()))
            except ImportError:
                logger.warning("7z support requires 'py7zr' package")
            except Exception as e:
                logger.warning("7z extraction failed: %s", e)

        return entries

    # ═══════════════════════════════════════════════════════════════════════════
    # PDF (hybrid: pymupdf text + VLM for scanned pages)
    # ═══════════════════════════════════════════════════════════════════════════

    async def _extract_pdf(self, file_bytes: bytes, max_tier: str) -> ExtractedDocument:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages: list[PageContent] = []
        method = "pymupdf"
        has_vlm_pages = False

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text = page.get_text("text").strip()
            image_list = page.get_images(full=True)
            has_images = len(image_list) > 0
            is_scan = len(page_text) < 100

            if (has_images or is_scan) and (is_scan or not page_text):
                try:
                    pix = page.get_pixmap(dpi=200)
                    image_bytes = pix.tobytes("png")
                    vlm_text = await self._call_vlm(image_bytes, max_tier)
                    page_content = PageContent(
                        page_number=page_num + 1, text=vlm_text,
                        images=[ImageDescription(description=vlm_text)])
                    has_vlm_pages = True
                except Exception as e:
                    logger.warning("VLM failed for PDF page %d: %s — using pymupdf text", page_num + 1, e)
                    page_content = PageContent(page_number=page_num + 1, text=page_text)
            else:
                page_content = PageContent(page_number=page_num + 1, text=page_text)
            pages.append(page_content)

        doc.close()
        if has_vlm_pages:
            method = "hybrid"

        full_text = "\n\n".join(f"--- Page {p.page_number} ---\n{p.text}" for p in pages if p.text)
        return ExtractedDocument(text=full_text, pages=pages, method=method,
                                 metadata={"page_count": len(pages), "has_images": has_vlm_pages})

    # ═══════════════════════════════════════════════════════════════════════════
    # Image (VLM)
    # ═══════════════════════════════════════════════════════════════════════════

    async def _extract_image(self, file_bytes: bytes, max_tier: str) -> ExtractedDocument:
        text = await self._call_vlm(file_bytes, max_tier)
        return ExtractedDocument(text=text, method="vlm", metadata={"type": "image"})

    async def _call_vlm(self, image_bytes: bytes, max_tier: str) -> str:
        """Extract text + visual description via router VLM (gRPC).

        Calls `RouterInferenceService.Generate` on `jervis-ollama-router:5501`.
        The router resolves tier + model via RequestContext (capability =
        VISUAL). `max_tier` is kept on the signature for call-site
        compatibility but isn't propagated — tier is client-based on the
        router side, surfaced here later when the extract API carries
        client_id.
        """
        import grpc.aio
        from jervis.common import enums_pb2, types_pb2
        from jervis.router import inference_pb2, inference_pb2_grpc
        from jervis_contracts.interceptors import prepare_context

        prompt = ("Extract all text, data, and describe all visual elements "
                  "(diagrams, charts, tables, images) from this document. "
                  "Preserve the original language. Output plain text.")

        target = self._router_grpc_target()
        async with grpc.aio.insecure_channel(
            target,
            options=[
                ("grpc.max_send_message_length", 32 * 1024 * 1024),
                ("grpc.max_receive_message_length", 32 * 1024 * 1024),
                ("grpc.keepalive_time_ms", 30_000),
                ("grpc.keepalive_timeout_ms", 10_000),
                ("grpc.keepalive_permit_without_calls", 1),
            ],
        ) as channel:
            stub = inference_pb2_grpc.RouterInferenceServiceStub(channel)
            ctx = types_pb2.RequestContext(
                scope=types_pb2.Scope(),
                priority=enums_pb2.PRIORITY_BACKGROUND,
                capability=enums_pb2.CAPABILITY_VISUAL,
                intent="document-extraction-vlm",
            )
            prepare_context(ctx)
            request = inference_pb2.GenerateRequest(
                ctx=ctx,
                model_hint="qwen3-vl-tool:latest",
                prompt=prompt,
                images=[image_bytes],
                options=inference_pb2.ChatOptions(temperature=0.0, num_predict=4096),
            )

            parts: list[str] = []
            async for chunk in stub.Generate(request):
                if chunk.response_delta:
                    parts.append(chunk.response_delta)

        text = "".join(parts)
        logger.info("VLM response: %d chars", len(text))
        return text

    def _router_grpc_target(self) -> str:
        """Strip scheme/port from configured router URL and target :5501."""
        u = self.router_url.rstrip("/")
        if "://" in u:
            u = u.split("://", 1)[1]
        host = u.split("/")[0].split(":")[0]
        return f"{host}:5501"
